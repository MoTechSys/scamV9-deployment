"""
خدمات الذكاء الاصطناعي - Manus API Proxy (OpenAI Compatible)
S-ACM - Smart Academic Content Management System

=== Phase 3: Manus Proxy Engine ===
1. SimpleKeyManager: Single key from .env (MANUS_API_KEY) - no DB rotation
2. SmartChunker: Uses AIConfiguration.chunk_size from DB
3. GeminiService: Uses OpenAI-compatible client via Manus API Proxy
4. All settings are Admin-editable via AIConfiguration

== Migration Notes ==
- Replaced google-genai with openai library
- HydraKeyManager simplified to SimpleKeyManager (single env key)
- base_url = https://api.manus.im/api/llm-proxy/v1
- All public interfaces preserved (GeminiService, QuestionMatrixConfig, etc.)
"""

from __future__ import annotations

import json
import hashlib
import logging
import os
import re
import time
import threading
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from enum import Enum
from functools import wraps
from pathlib import Path
from typing import Optional, List, Dict, Any, Callable, TypeVar
from datetime import datetime

from django.conf import settings
from django.core.cache import cache
from django.utils import timezone

# ========== Logging ==========
logger = logging.getLogger('ai_features')


# ========== Constants (Fallbacks - DB config takes priority) ==========
FALLBACK_MODEL = getattr(settings, 'AI_MODEL_NAME', None) or os.getenv('AI_MODEL_NAME', 'gpt-4.1-mini')
MANUS_BASE_URL = getattr(settings, 'MANUS_BASE_URL', None) or os.getenv('MANUS_BASE_URL', 'https://api.manus.im/api/llm-proxy/v1')
FALLBACK_CHUNK_SIZE = 30000
FALLBACK_CHUNK_OVERLAP = 500
FALLBACK_MAX_OUTPUT_TOKENS = 2000
FALLBACK_TEMPERATURE = 0.3
CACHE_TIMEOUT = 3600
MAX_RETRIES = 3
AI_OUTPUT_DIR = 'ai_generated'

# Legacy compatibility aliases
GEMINI_MODEL = FALLBACK_MODEL
MAX_INPUT_LENGTH = FALLBACK_CHUNK_SIZE
CHUNK_SIZE = 8000
CHUNK_OVERLAP = 500


# ========== Custom Exceptions ==========

class GeminiError(Exception):
    """Base exception for AI-related errors (kept for backward compat)."""
    pass


class GeminiConfigurationError(GeminiError):
    """Raised when AI is not properly configured."""
    pass


class GeminiAPIError(GeminiError):
    """Raised when AI API returns an error."""
    def __init__(self, message: str, status_code: Optional[int] = None):
        super().__init__(message)
        self.status_code = status_code


class GeminiRateLimitError(GeminiAPIError):
    """Raised when rate limit is exceeded."""
    pass


class GeminiServiceDisabledError(GeminiError):
    """Raised when AI service is disabled by admin."""
    pass


class TextExtractionError(GeminiError):
    """Raised when text extraction from file fails."""
    pass


# ========== Enums ==========

class QuestionType(Enum):
    MCQ = "mcq"
    TRUE_FALSE = "true_false"
    SHORT_ANSWER = "short_answer"
    MIXED = "mixed"


class ContentType(Enum):
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    TEXT = "text"
    UNKNOWN = "unknown"


# ========== Data Classes ==========

@dataclass
class Question:
    """نموذج سؤال واحد."""
    type: str
    question: str
    answer: str
    options: Optional[List[str]] = None
    explanation: Optional[str] = None
    score: float = 1.0

    def to_dict(self) -> Dict[str, Any]:
        result = {
            'type': self.type,
            'question': self.question,
            'answer': self.answer,
            'score': self.score,
        }
        if self.options:
            result['options'] = self.options
        if self.explanation:
            result['explanation'] = self.explanation
        return result


@dataclass
class QuestionMatrixConfig:
    """تكوين مصفوفة الأسئلة."""
    mcq_count: int = 0
    mcq_score: float = 2.0
    true_false_count: int = 0
    true_false_score: float = 1.0
    short_answer_count: int = 0
    short_answer_score: float = 3.0

    @property
    def total_questions(self) -> int:
        return self.mcq_count + self.true_false_count + self.short_answer_count

    @property
    def total_score(self) -> float:
        return (
            self.mcq_count * self.mcq_score +
            self.true_false_count * self.true_false_score +
            self.short_answer_count * self.short_answer_score
        )

    def to_dict(self) -> Dict[str, Any]:
        return {
            'mcq_count': self.mcq_count, 'mcq_score': self.mcq_score,
            'true_false_count': self.true_false_count, 'true_false_score': self.true_false_score,
            'short_answer_count': self.short_answer_count, 'short_answer_score': self.short_answer_score,
            'total_questions': self.total_questions, 'total_score': self.total_score,
        }


@dataclass
class AIResponse:
    """نموذج استجابة AI."""
    success: bool
    data: Any = None
    error: Optional[str] = None
    cached: bool = False
    md_file_path: Optional[str] = None


# ========================================================================
# Simple Key Manager (Manus API - Single Key from .env)
# ========================================================================

class HydraKeyManager:
    """
    Simplified Key Manager - Single Manus API Key from environment.

    Replaces the old DB-based Round-Robin system with a simple
    environment variable approach for Manus API Proxy.

    Usage:
        manager = HydraKeyManager()
        key = manager.get_api_key()
    """
    _instance = None
    _lock = threading.Lock()

    def __new__(cls):
        if cls._instance is None:
            with cls._lock:
                if cls._instance is None:
                    cls._instance = super().__new__(cls)
                    cls._instance._initialized = False
        return cls._instance

    def __init__(self):
        if self._initialized:
            return
        self._api_key = None
        self._load_key()
        self._initialized = True

    def _load_key(self):
        """Load the Manus API key from environment / settings."""
        self._api_key = (
            getattr(settings, 'MANUS_API_KEY', '') or
            os.getenv('MANUS_API_KEY', '')
        )
        if self._api_key:
            logger.info("HydraKeyManager: Manus API key loaded from environment")
        else:
            logger.warning("HydraKeyManager: No MANUS_API_KEY found in environment!")

    def get_api_key(self) -> str:
        """Get the Manus API key."""
        if not self._api_key:
            raise GeminiConfigurationError(
                "مفتاح MANUS_API_KEY غير موجود. أضف المفتاح في ملف .env"
            )
        return self._api_key

    def get_next_key(self):
        """
        Backward-compatible interface.
        Returns (None, api_key) to match old signature.
        """
        return None, self.get_api_key()

    @property
    def total_keys(self) -> int:
        return 1 if self._api_key else 0

    @property
    def has_keys(self) -> bool:
        return bool(self._api_key)

    def rotate_after_error(self, failed_key_obj=None, error_msg: str = '', is_rate_limit: bool = False):
        """No-op for single key system. Kept for backward compat."""
        logger.warning(f"HydraKeyManager: Error occurred: {error_msg[:100]}")

    def get_health_status(self) -> List[Dict[str, Any]]:
        """Get health status (simplified for single key)."""
        if self._api_key:
            return [{
                'id': 0,
                'label': 'Manus API Key (ENV)',
                'hint': self._api_key[-4:] if len(self._api_key) >= 4 else '****',
                'status': 'active',
                'is_available': True,
                'error_count': 0,
                'total_requests': 0,
                'last_latency_ms': 0,
                'last_success_at': None,
                'last_error': None,
                'rpm_limit': 0,
                'tokens_used_today': 0,
                'cooldown_until': None,
            }]
        return []


# Legacy compatibility: Alias
APIKeyManager = HydraKeyManager


# ========================================================================
# Smart Text Chunking (DB-Configured)
# ========================================================================

class SmartChunker:
    """
    تقسيم ذكي للنصوص الكبيرة.
    Reads chunk_size and overlap from AIConfiguration (DB).
    """

    def __init__(self, chunk_size: int = None, overlap: int = None):
        if chunk_size is None or overlap is None:
            try:
                from apps.ai_features.models import AIConfiguration
                config = AIConfiguration.get_config()
                self.chunk_size = chunk_size or config.chunk_size
                self.overlap = overlap or config.chunk_overlap
            except Exception:
                self.chunk_size = chunk_size or FALLBACK_CHUNK_SIZE
                self.overlap = overlap or FALLBACK_CHUNK_OVERLAP
        else:
            self.chunk_size = chunk_size
            self.overlap = overlap

    def chunk_text(self, text: str) -> List[str]:
        """تقسيم النص إلى أجزاء ذكية."""
        if not text:
            return []
        if len(text) <= self.chunk_size:
            return [text]

        chunks = []
        paragraphs = text.split('\n\n')
        current_chunk = ""

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            if len(para) > self.chunk_size:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    current_chunk = ""
                sentence_chunks = self._split_by_sentences(para)
                chunks.extend(sentence_chunks)
                continue

            if len(current_chunk) + len(para) + 2 > self.chunk_size:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    overlap_text = current_chunk[-self.overlap:] if len(current_chunk) > self.overlap else ""
                    current_chunk = overlap_text + "\n\n" + para
                else:
                    current_chunk = para
            else:
                current_chunk += ("\n\n" if current_chunk else "") + para

        if current_chunk.strip():
            chunks.append(current_chunk.strip())

        logger.info(f"SmartChunker: Split {len(text)} chars into {len(chunks)} chunks (size={self.chunk_size})")
        return chunks

    def _split_by_sentences(self, text: str) -> List[str]:
        """تقسيم النص بالجمل."""
        separators = ['. ', '.\n', '\u3002', '\u061f ', '? ', '! ', '\uff01 ', '.\t']
        sentences = [text]
        for sep in separators:
            new_sentences = []
            for s in sentences:
                parts = s.split(sep)
                for i, part in enumerate(parts):
                    if i < len(parts) - 1:
                        new_sentences.append(part + sep.strip())
                    else:
                        if part.strip():
                            new_sentences.append(part)
            sentences = new_sentences

        chunks = []
        current = ""
        for sentence in sentences:
            if len(current) + len(sentence) + 1 > self.chunk_size:
                if current:
                    chunks.append(current.strip())
                current = sentence
            else:
                current += " " + sentence if current else sentence

        if current.strip():
            chunks.append(current.strip())

        return chunks


# ========== File-Based AI Storage ==========

class AIFileStorage:
    """مخزن ملفات AI - يحفظ المخرجات كملفات .md في media/ai_generated/."""

    def __init__(self):
        self.base_dir = Path(settings.MEDIA_ROOT) / AI_OUTPUT_DIR
        self.base_dir.mkdir(parents=True, exist_ok=True)

    def save_summary(self, file_id: int, content: str, metadata: Optional[Dict] = None) -> str:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"summary_{file_id}_{timestamp}.md"
        return self._save_file(filename, content, metadata, "summary")

    def save_questions(self, file_id: int, questions_data: List[Dict], metadata: Optional[Dict] = None) -> str:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"questions_{file_id}_{timestamp}.md"
        md_content = self._questions_to_markdown(questions_data, metadata)
        return self._save_file(filename, md_content, None, "questions")

    def save_chat_answer(self, file_id: int, user_id: int, question: str, answer: str) -> str:
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"chat_{file_id}_{user_id}_{timestamp}.md"
        content = f"# \u0633\u0624\u0627\u0644\n\n{question}\n\n# \u0625\u062c\u0627\u0628\u0629\n\n{answer}\n"
        return self._save_file(filename, content, None, "chat")

    def read_file(self, relative_path: str) -> Optional[str]:
        full_path = Path(settings.MEDIA_ROOT) / relative_path
        try:
            if full_path.exists():
                return full_path.read_text(encoding='utf-8')
        except Exception as e:
            logger.error(f"AIFileStorage: Failed to read {full_path}: {e}")
        return None

    def delete_file(self, relative_path: str) -> bool:
        full_path = Path(settings.MEDIA_ROOT) / relative_path
        try:
            if full_path.exists():
                full_path.unlink()
                logger.info(f"AIFileStorage: Deleted {relative_path}")
                return True
        except Exception as e:
            logger.error(f"AIFileStorage: Failed to delete {relative_path}: {e}")
        return False

    def _save_file(self, filename: str, content: str, metadata: Optional[Dict], category: str) -> str:
        category_dir = self.base_dir / category
        category_dir.mkdir(parents=True, exist_ok=True)

        filepath = category_dir / filename
        header = ""
        if metadata:
            header = "---\n"
            for k, v in metadata.items():
                header += f"{k}: {v}\n"
            header += "---\n\n"

        filepath.write_text(header + content, encoding='utf-8')
        relative_path = str(Path(AI_OUTPUT_DIR) / category / filename)
        logger.info(f"AIFileStorage: Saved {relative_path} ({len(content)} chars)")
        return relative_path

    def _questions_to_markdown(self, questions: List[Dict], metadata: Optional[Dict] = None) -> str:
        lines = ["# \u0628\u0646\u0643 \u0627\u0644\u0623\u0633\u0626\u0644\u0629 \u0627\u0644\u0645\u064f\u0648\u0644\u064e\u0651\u062f\u0629 \u0628\u0627\u0644\u0630\u0643\u0627\u0621 \u0627\u0644\u0627\u0635\u0637\u0646\u0627\u0639\u064a\n"]

        if metadata:
            lines.append(f"**\u0627\u0644\u0645\u0635\u062f\u0631:** {metadata.get('source_file', '\u063a\u064a\u0631 \u0645\u062d\u062f\u062f')}")
            lines.append(f"**\u0627\u0644\u062a\u0627\u0631\u064a\u062e:** {metadata.get('date', datetime.now().strftime('%Y-%m-%d'))}")
            lines.append(f"**\u0625\u062c\u0645\u0627\u0644\u064a \u0627\u0644\u062f\u0631\u062c\u0627\u062a:** {metadata.get('total_score', '-')}")
            lines.append("")

        type_labels = {
            'mcq': '\u0627\u062e\u062a\u064a\u0627\u0631 \u0645\u0646 \u0645\u062a\u0639\u062f\u062f',
            'true_false': '\u0635\u062d \u0648\u062e\u0637\u0623',
            'short_answer': '\u0625\u062c\u0627\u0628\u0629 \u0642\u0635\u064a\u0631\u0629'
        }

        for i, q in enumerate(questions, 1):
            q_type = q.get('type', 'short_answer')
            score = q.get('score', 1.0)
            label = type_labels.get(q_type, q_type)

            lines.append(f"## \u0627\u0644\u0633\u0624\u0627\u0644 {i} ({label}) - [{score} \u062f\u0631\u062c\u0629]")
            lines.append(f"\n{q.get('question', '')}\n")

            options = q.get('options')
            if options and isinstance(options, list):
                for j, opt in enumerate(options):
                    letter = chr(ord('\u0623') + j) if j < 4 else chr(ord('a') + j)
                    lines.append(f"- {letter}) {opt}")
                lines.append("")

            lines.append(f"**\u0627\u0644\u0625\u062c\u0627\u0628\u0629:** {q.get('answer', '')}")

            explanation = q.get('explanation')
            if explanation:
                lines.append(f"\n**\u0627\u0644\u0634\u0631\u062d:** {explanation}")
            lines.append("\n---\n")

        return "\n".join(lines)


# ========== Decorators ==========

T = TypeVar('T')


def cache_result(timeout: int = CACHE_TIMEOUT):
    """Decorator لتخزين نتائج AI في الكاش."""
    def decorator(func: Callable[..., T]) -> Callable[..., T]:
        @wraps(func)
        def wrapper(self, text: str, *args, **kwargs) -> T:
            cache_key = _generate_cache_key(func.__name__, text, args, kwargs)
            cached_result = cache.get(cache_key)
            if cached_result is not None:
                logger.debug(f"Cache hit for {func.__name__}")
                return cached_result
            result = func(self, text, *args, **kwargs)
            if result is not None:
                cache.set(cache_key, result, timeout)
            return result
        return wrapper
    return decorator


def _generate_cache_key(func_name: str, text: str, args: tuple, kwargs: dict) -> str:
    content = f"{func_name}:{text[:200]}:{str(args)}:{str(sorted(kwargs.items()))}"
    return f"ai:{hashlib.md5(content.encode()).hexdigest()}"


# ========== Text Extractors ==========

class TextExtractor(ABC):
    @abstractmethod
    def extract(self, file_path: Path) -> str:
        pass

    @abstractmethod
    def supports(self, file_path: Path) -> bool:
        pass


class PDFExtractor(TextExtractor):
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == '.pdf'

    def extract(self, file_path: Path) -> str:
        try:
            import pdfplumber
        except ImportError:
            raise TextExtractionError("pdfplumber not installed. Run: pip install pdfplumber")
        try:
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            return "\n".join(text_parts)
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from PDF: {e}")


class DocxExtractor(TextExtractor):
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in ['.docx', '.doc']

    def extract(self, file_path: Path) -> str:
        try:
            from docx import Document
        except ImportError:
            raise TextExtractionError("python-docx not installed. Run: pip install python-docx")
        try:
            doc = Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs if para.text)
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from DOCX: {e}")


class PptxExtractor(TextExtractor):
    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == '.pptx'

    def extract(self, file_path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise TextExtractionError("python-pptx not installed. Run: pip install python-pptx")
        try:
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as e:
            raise TextExtractionError(f"Failed to extract text from PPTX: {e}")


class PlainTextExtractor(TextExtractor):
    SUPPORTED_EXTENSIONS = {'.txt', '.md', '.rst', '.csv'}

    def supports(self, file_path: Path) -> bool:
        return file_path.suffix.lower() in self.SUPPORTED_EXTENSIONS

    def extract(self, file_path: Path) -> str:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            with open(file_path, 'r', encoding='cp1256') as f:
                return f.read()
        except Exception as e:
            raise TextExtractionError(f"Failed to read text file: {e}")


class TextExtractorFactory:
    _extractors: List[TextExtractor] = [
        PDFExtractor(),
        DocxExtractor(),
        PptxExtractor(),
        PlainTextExtractor(),
    ]

    @classmethod
    def get_extractor(cls, file_path: Path) -> Optional[TextExtractor]:
        for extractor in cls._extractors:
            if extractor.supports(file_path):
                return extractor
        return None

    @classmethod
    def extract_text(cls, file_path: Path) -> str:
        extractor = cls.get_extractor(file_path)
        if extractor is None:
            raise TextExtractionError(f"Unsupported file type: {file_path.suffix}")
        return extractor.extract(file_path)


# ========================================================================
# AI Service (Manus API Proxy - OpenAI Compatible)
# ========================================================================

def _get_ai_config():
    """Helper to get AI configuration from DB with fallback."""
    try:
        from apps.ai_features.models import AIConfiguration
        return AIConfiguration.get_config()
    except Exception:
        return None


class GeminiService:
    """
    خدمة الذكاء الاصطناعي عبر Manus API Proxy (OpenAI-compatible).

    === Manus Proxy Engine ===
    - Uses OpenAI Python client with Manus base_url
    - Model, tokens, temperature: Read from AIConfiguration (DB)
    - API Key: Single MANUS_API_KEY from environment
    - Chunk size: Read from AIConfiguration (DB)
    - Service toggle: Can be disabled from Admin panel

    === Admin Editable (No Code Touch) ===
    - Change model: Admin -> AI Configuration -> active_model
    - Disable service: Admin -> AI Configuration -> is_service_enabled
    """

    def __init__(self, model: str = None):
        config = _get_ai_config()
        self._model_name = model or (config.active_model if config else FALLBACK_MODEL)
        self._key_manager = HydraKeyManager()
        self._chunker = SmartChunker()
        self._storage = AIFileStorage()
        self._client = None
        self._initialize_client()

    def _check_service_enabled(self):
        """Check if AI service is enabled by admin."""
        config = _get_ai_config()
        if config and not config.is_service_enabled:
            msg = config.maintenance_message or '\u062e\u062f\u0645\u0629 \u0627\u0644\u0630\u0643\u0627\u0621 \u0627\u0644\u0627\u0635\u0637\u0646\u0627\u0639\u064a \u0645\u062a\u0648\u0642\u0641\u0629 \u0645\u0624\u0642\u062a\u0627\u064b.'
            raise GeminiServiceDisabledError(msg)

    def _initialize_client(self) -> None:
        """تهيئة عميل OpenAI عبر Manus API Proxy."""
        if not self._key_manager.has_keys:
            logger.warning("GeminiService: No MANUS_API_KEY available. Service will be limited.")
            return
        try:
            from openai import OpenAI

            api_key = self._key_manager.get_api_key()
            self._client = OpenAI(
                api_key=api_key,
                base_url=MANUS_BASE_URL,
            )
            logger.info(f"GeminiService initialized with Manus Proxy | model: {self._model_name} | base_url: {MANUS_BASE_URL}")
        except ImportError:
            raise GeminiConfigurationError("openai not installed. Run: pip install openai")
        except GeminiConfigurationError:
            logger.warning("GeminiService: No API key available for initialization.")
        except Exception as e:
            raise GeminiConfigurationError(f"Failed to initialize OpenAI client: {e}")

    @property
    def is_available(self) -> bool:
        return self._client is not None

    @property
    def storage(self) -> AIFileStorage:
        return self._storage

    def _generate_content(self, prompt: str, max_tokens: int = None) -> str:
        """توليد محتوى عبر Manus API Proxy (OpenAI-compatible)."""
        self._check_service_enabled()

        if not self.is_available:
            raise GeminiConfigurationError("AI client not initialized. Check MANUS_API_KEY.")

        # Get config from DB
        config = _get_ai_config()
        if max_tokens is None:
            max_tokens = config.max_output_tokens if config else FALLBACK_MAX_OUTPUT_TOKENS
        temperature = config.temperature if config else FALLBACK_TEMPERATURE

        last_exception = None

        for attempt in range(MAX_RETRIES):
            try:
                start_ms = int(time.time() * 1000)

                # Use OpenAI chat.completions.create format
                response = self._client.chat.completions.create(
                    model=self._model_name,
                    messages=[
                        {"role": "user", "content": prompt}
                    ],
                    max_tokens=max_tokens,
                    temperature=temperature,
                )

                latency_ms = int(time.time() * 1000) - start_ms

                # Extract response text
                if response.choices and response.choices[0].message.content:
                    result_text = response.choices[0].message.content.strip()
                    logger.info(f"AI response received in {latency_ms}ms (model={self._model_name})")
                    return result_text
                else:
                    raise GeminiAPIError("Empty response from AI API")

            except GeminiServiceDisabledError:
                raise
            except GeminiError:
                raise
            except Exception as e:
                last_exception = e
                error_str = str(e).lower()

                is_rate_limit = any(kw in error_str for kw in [
                    "rate", "quota", "429", "resource_exhausted", "too many requests"
                ])

                if is_rate_limit:
                    logger.warning(f"Rate limit on attempt {attempt + 1}, waiting...")
                    wait_time = min(5.0 * (2 ** attempt), 30)
                    time.sleep(wait_time)
                    continue
                elif "invalid" in error_str and ("key" in error_str or "auth" in error_str):
                    logger.error(f"Invalid API key or auth error: {e}")
                    raise GeminiConfigurationError(f"Authentication error: {e}")
                else:
                    if attempt < MAX_RETRIES - 1:
                        wait_time = min(2.0 * (2 ** attempt), 15)
                        logger.warning(f"API error on attempt {attempt + 1}, retrying in {wait_time}s: {e}")
                        time.sleep(wait_time)
                        continue
                    raise GeminiAPIError(f"AI API error: {e}")

        if last_exception and any(kw in str(last_exception).lower() for kw in ["quota", "429", "rate"]):
            raise GeminiRateLimitError(
                "\u23f3 \u062a\u0645 \u062a\u062c\u0627\u0648\u0632 \u0627\u0644\u062d\u062f \u0627\u0644\u0645\u0633\u0645\u0648\u062d \u0644\u0640 API. \u064a\u0631\u062c\u0649 \u0627\u0644\u0627\u0646\u062a\u0638\u0627\u0631 \u062f\u0642\u064a\u0642\u0629 \u062b\u0645 \u0627\u0644\u0645\u062d\u0627\u0648\u0644\u0629 \u0645\u0631\u0629 \u0623\u062e\u0631\u0649."
            )
        raise last_exception or GeminiAPIError("AI API request failed after all retries")

    # ========== Public Methods ==========

    def extract_text_from_file(self, file_obj) -> Optional[str]:
        """استخراج النص من كائن LectureFile."""
        if not file_obj.local_file:
            logger.warning(f"File {file_obj.id} has no local file")
            return None
        try:
            file_path = Path(file_obj.local_file.path)
            text = TextExtractorFactory.extract_text(file_path)
            logger.info(f"Extracted {len(text)} characters from {file_path.name}")
            return text
        except TextExtractionError as e:
            logger.error(f"Text extraction failed for file {file_obj.id}: {e}")
            return None

    @cache_result(timeout=CACHE_TIMEOUT)
    def generate_summary(self, text: str, max_length: int = 500, user_notes: str = "") -> str:
        """توليد تلخيص للنص مع دعم Smart Chunking."""
        chunks = self._chunker.chunk_text(text)

        if len(chunks) <= 1:
            return self._generate_single_summary(text, max_length, user_notes)

        chunk_summaries = []
        for i, chunk in enumerate(chunks):
            try:
                summary = self._generate_single_summary(
                    chunk, max_length=200,
                    user_notes=f"\u0647\u0630\u0627 \u0627\u0644\u062c\u0632\u0621 {i+1} \u0645\u0646 {len(chunks)}. {user_notes}"
                )
                chunk_summaries.append(summary)
            except GeminiError:
                continue

        if not chunk_summaries:
            return self._fallback_summary(text, max_length)

        combined = "\n\n".join(chunk_summaries)
        return self._generate_single_summary(
            combined, max_length,
            user_notes=f"\u0647\u0630\u0627 \u0645\u0644\u062e\u0635 \u0645\u064f\u062c\u0645\u0651\u0639 \u0645\u0646 {len(chunks)} \u0623\u062c\u0632\u0627\u0621. \u0623\u0639\u062f \u0635\u064a\u0627\u063a\u062a\u0647 \u0643\u0645\u0644\u062e\u0635 \u0645\u062a\u0645\u0627\u0633\u0643 \u0648\u0627\u062d\u062f. {user_notes}"
        )

    def _generate_single_summary(self, text: str, max_length: int = 500, user_notes: str = "") -> str:
        notes_section = f"\n\nUSER_INSTRUCTION: {user_notes}" if user_notes else ""
        config = _get_ai_config()
        input_limit = config.chunk_size if config else FALLBACK_CHUNK_SIZE

        prompt = f"""ROLE: \u0623\u0646\u062a \u0645\u0633\u0627\u0639\u062f \u0623\u0643\u0627\u062f\u064a\u0645\u064a \u062e\u0628\u064a\u0631 \u0645\u062a\u062e\u0635\u0635 \u0641\u064a \u062a\u0644\u062e\u064a\u0635 \u0627\u0644\u0645\u062d\u062a\u0648\u0649 \u0627\u0644\u062a\u0639\u0644\u064a\u0645\u064a.

TASK: \u062a\u0644\u062e\u064a\u0635
\u0642\u0645 \u0628\u062a\u0644\u062e\u064a\u0635 \u0627\u0644\u0646\u0635 \u0627\u0644\u062a\u0627\u0644\u064a \u0628\u0634\u0643\u0644 \u0645\u062e\u062a\u0635\u0631 \u0648\u0645\u0641\u064a\u062f \u0628\u0635\u064a\u063a\u0629 Markdown. \u0631\u0643\u0632 \u0639\u0644\u0649:
- \u0627\u0644\u0646\u0642\u0627\u0637 \u0627\u0644\u0631\u0626\u064a\u0633\u064a\u0629 \u0648\u0627\u0644\u0645\u0641\u0627\u0647\u064a\u0645 \u0627\u0644\u0623\u0633\u0627\u0633\u064a\u0629
- \u0627\u0644\u0645\u0639\u0644\u0648\u0645\u0627\u062a \u0627\u0644\u0623\u0643\u062b\u0631 \u0623\u0647\u0645\u064a\u0629
- \u0627\u0644\u062d\u0641\u0627\u0638 \u0639\u0644\u0649 \u0627\u0644\u062f\u0642\u0629 \u0627\u0644\u0639\u0644\u0645\u064a\u0629
- \u0627\u0633\u062a\u062e\u062f\u0627\u0645 \u0639\u0646\u0627\u0648\u064a\u0646 \u0648\u0642\u0648\u0627\u0626\u0645 \u0644\u062a\u0646\u0638\u064a\u0645 \u0627\u0644\u0645\u062d\u062a\u0648\u0649
{notes_section}

OUTPUT_FORMAT: Markdown (\u0645\u0639 \u0639\u0646\u0627\u0648\u064a\u0646 \u0648\u0642\u0648\u0627\u0626\u0645 \u0648\u062a\u0646\u0633\u064a\u0642 \u0648\u0627\u0636\u062d)
LANGUAGE: Arabic (\u0645\u0627 \u0644\u0645 \u064a\u064f\u062d\u062f\u062f \u062e\u0644\u0627\u0641 \u0630\u0644\u0643)

CONTEXT:
{text[:input_limit]}

\u0627\u0644\u062a\u0644\u062e\u064a\u0635 (\u0628\u062d\u062f \u0623\u0642\u0635\u0649 {max_length} \u0643\u0644\u0645\u0629\u060c \u0628\u0635\u064a\u063a\u0629 Markdown):"""

        try:
            return self._generate_content(prompt, max_tokens=max_length * 3)
        except GeminiError as e:
            logger.error(f"Summary generation failed: {e}")
            return self._fallback_summary(text, max_length)

    def _fallback_summary(self, text: str, max_length: int) -> str:
        sentences = text.replace('\n', ' ').split('.')
        summary = ""
        for sentence in sentences:
            sentence = sentence.strip()
            if sentence and len(summary) + len(sentence) < max_length * 4:
                summary += sentence + ". "
            elif len(summary) > 100:
                break
        return summary.strip() or text[:max_length * 4] + "..."

    def generate_questions_matrix(
        self, text: str, matrix: QuestionMatrixConfig, user_notes: str = ""
    ) -> List[Dict[str, Any]]:
        """توليد أسئلة حسب المصفوفة المحددة."""
        if matrix.total_questions == 0:
            return []

        chunks = self._chunker.chunk_text(text)
        config = _get_ai_config()
        input_limit = config.chunk_size if config else FALLBACK_CHUNK_SIZE

        source_text = chunks[0] if chunks else text[:input_limit]
        if len(chunks) > 1:
            source_text = "\n\n---\n\n".join(chunks[:3])

        notes_section = f"\n\nUSER_INSTRUCTION: {user_notes}" if user_notes else ""

        parts = []
        if matrix.mcq_count > 0:
            parts.append(f"- {matrix.mcq_count} \u0633\u0624\u0627\u0644 \u0627\u062e\u062a\u064a\u0627\u0631 \u0645\u0646 \u0645\u062a\u0639\u062f\u062f (mcq) - \u0643\u0644 \u0633\u0624\u0627\u0644 {matrix.mcq_score} \u062f\u0631\u062c\u0629")
        if matrix.true_false_count > 0:
            parts.append(f"- {matrix.true_false_count} \u0633\u0624\u0627\u0644 \u0635\u062d \u0648\u062e\u0637\u0623 (true_false) - \u0643\u0644 \u0633\u0624\u0627\u0644 {matrix.true_false_score} \u062f\u0631\u062c\u0629")
        if matrix.short_answer_count > 0:
            parts.append(f"- {matrix.short_answer_count} \u0633\u0624\u0627\u0644 \u0625\u062c\u0627\u0628\u0629 \u0642\u0635\u064a\u0631\u0629 (short_answer) - \u0643\u0644 \u0633\u0624\u0627\u0644 {matrix.short_answer_score} \u062f\u0631\u062c\u0629")

        matrix_text = "\n".join(parts)

        prompt = f"""ROLE: \u0623\u0646\u062a \u0645\u062f\u0631\u0633 \u062c\u0627\u0645\u0639\u064a \u062e\u0628\u064a\u0631 \u0645\u062a\u062e\u0635\u0635 \u0641\u064a \u0625\u0646\u0634\u0627\u0621 \u0623\u0633\u0626\u0644\u0629 \u0627\u062e\u062a\u0628\u0627\u0631\u064a\u0629 \u0623\u0643\u0627\u062f\u064a\u0645\u064a\u0629.

TASK: \u062a\u0648\u0644\u064a\u062f \u0623\u0633\u0626\u0644\u0629 \u0627\u062e\u062a\u0628\u0627\u0631\u064a\u0629
CONFIG:
{matrix_text}

\u0625\u062c\u0645\u0627\u0644\u064a: {matrix.total_questions} \u0633\u0624\u0627\u0644 | \u0627\u0644\u062f\u0631\u062c\u0629 \u0627\u0644\u0643\u0644\u064a\u0629: {matrix.total_score}
{notes_section}

\u0623\u0631\u062c\u0639 \u0627\u0644\u0625\u062c\u0627\u0628\u0629 \u0628\u0635\u064a\u063a\u0629 JSON \u0641\u0642\u0637 \u0628\u062f\u0648\u0646 \u0623\u064a \u0646\u0635 \u0625\u0636\u0627\u0641\u064a\u060c \u0643\u0645\u0635\u0641\u0648\u0641\u0629:
[
    {{
        "type": "mcq" \u0623\u0648 "true_false" \u0623\u0648 "short_answer",
        "question": "\u0646\u0635 \u0627\u0644\u0633\u0624\u0627\u0644",
        "options": ["\u062e\u064a\u0627\u06311", "\u062e\u064a\u0627\u06312", "\u062e\u064a\u0627\u06313", "\u062e\u064a\u0627\u06314"],
        "answer": "\u0627\u0644\u0625\u062c\u0627\u0628\u0629 \u0627\u0644\u0635\u062d\u064a\u062d\u0629",
        "explanation": "\u0634\u0631\u062d \u0645\u062e\u062a\u0635\u0631",
        "score": \u0627\u0644\u062f\u0631\u062c\u0629_\u0643\u0631\u0642\u0645
    }}
]

\u0645\u0644\u0627\u062d\u0638\u0627\u062a:
- \u0644\u0644\u0623\u0633\u0626\u0644\u0629 true_false: \u0627\u0644\u062e\u064a\u0627\u0631\u0627\u062a ["\u0635\u062d", "\u062e\u0637\u0623"] \u0641\u0642\u0637
- \u0644\u0644\u0623\u0633\u0626\u0644\u0629 short_answer: \u0644\u0627 \u062a\u0636\u0639 options (\u0627\u062c\u0639\u0644\u0647\u0627 null)
- \u062a\u0623\u0643\u062f \u0623\u0646 \u0627\u0644\u0623\u0633\u0626\u0644\u0629 \u0645\u062a\u0646\u0648\u0639\u0629 \u0648\u062a\u063a\u0637\u064a \u0623\u062c\u0632\u0627\u0621 \u0645\u062e\u062a\u0644\u0641\u0629 \u0645\u0646 \u0627\u0644\u0646\u0635

\u0627\u0644\u0646\u0635:
{source_text[:input_limit]}

\u0627\u0644\u0623\u0633\u0626\u0644\u0629 (JSON \u0641\u0642\u0637):"""

        try:
            result = self._generate_content(prompt, max_tokens=4000)
            return self._parse_questions_json(result)
        except GeminiError as e:
            logger.error(f"Question matrix generation failed: {e}")
            return self._fallback_questions(matrix.total_questions)

    @cache_result(timeout=CACHE_TIMEOUT)
    def generate_questions(
        self, text: str, question_type: QuestionType = QuestionType.MIXED,
        num_questions: int = 5, user_notes: str = ""
    ) -> List[Dict[str, Any]]:
        """توليد أسئلة (واجهة متوافقة مع الإصدار القديم)."""
        matrix = QuestionMatrixConfig()
        if question_type == QuestionType.MCQ:
            matrix.mcq_count = num_questions
        elif question_type == QuestionType.TRUE_FALSE:
            matrix.true_false_count = num_questions
        elif question_type == QuestionType.SHORT_ANSWER:
            matrix.short_answer_count = num_questions
        else:
            matrix.mcq_count = max(1, num_questions // 3)
            matrix.true_false_count = max(1, num_questions // 3)
            matrix.short_answer_count = num_questions - matrix.mcq_count - matrix.true_false_count
        return self.generate_questions_matrix(text, matrix, user_notes)

    def _parse_questions_json(self, result: str) -> List[Dict[str, Any]]:
        result = result.strip()
        if '```json' in result:
            result = result.split('```json')[1].split('```')[0]
        elif '```' in result:
            result = result.split('```')[1].split('```')[0]
        result = result.strip()

        try:
            questions = json.loads(result)
            if isinstance(questions, list):
                return questions
            return []
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse questions JSON: {e}")
            return []

    def _fallback_questions(self, num_questions: int) -> List[Dict[str, Any]]:
        return [{
            'type': 'short_answer',
            'question': '\u0645\u0627 \u0647\u064a \u0627\u0644\u0641\u0643\u0631\u0629 \u0627\u0644\u0631\u0626\u064a\u0633\u064a\u0629 \u0641\u064a \u0647\u0630\u0627 \u0627\u0644\u0646\u0635\u061f',
            'answer': '\u0631\u0627\u062c\u0639 \u0627\u0644\u0646\u0635 \u0644\u0644\u0625\u062c\u0627\u0628\u0629',
            'explanation': '\u0633\u0624\u0627\u0644 \u062a\u0644\u0642\u0627\u0626\u064a - \u062e\u062f\u0645\u0629 AI \u063a\u064a\u0631 \u0645\u062a\u0627\u062d\u0629 \u062d\u0627\u0644\u064a\u0627\u064b',
            'score': 1.0
        }]

    def ask_document(self, text: str, question: str, user_notes: str = "") -> str:
        """الإجابة على سؤال من سياق المستند."""
        chunks = self._chunker.chunk_text(text)
        config = _get_ai_config()
        input_limit = config.chunk_size if config else FALLBACK_CHUNK_SIZE
        context = chunks[0] if chunks else text[:input_limit]

        if len(chunks) > 1:
            context = self._find_relevant_chunks(chunks, question)

        notes_section = f"\n\nUSER_INSTRUCTION: {user_notes}" if user_notes else ""

        prompt = f"""ROLE: \u0623\u0646\u062a \u0645\u0633\u0627\u0639\u062f \u0623\u0643\u0627\u062f\u064a\u0645\u064a \u062e\u0628\u064a\u0631 \u064a\u062c\u064a\u0628 \u0639\u0644\u0649 \u0627\u0644\u0623\u0633\u0626\u0644\u0629 \u0628\u0646\u0627\u0621\u064b \u0639\u0644\u0649 \u0645\u062d\u062a\u0648\u0649 \u0627\u0644\u0645\u0633\u062a\u0646\u062f\u0627\u062a.

TASK: \u0627\u0644\u0625\u062c\u0627\u0628\u0629 \u0639\u0644\u0649 \u0633\u0624\u0627\u0644 \u0623\u0643\u0627\u062f\u064a\u0645\u064a
\u0642\u0648\u0627\u0639\u062f:
1. \u0623\u062c\u0628 \u0628\u0646\u0627\u0621\u064b \u0639\u0644\u0649 \u0627\u0644\u0645\u062d\u062a\u0648\u0649 \u0627\u0644\u0645\u0642\u062f\u0645 \u0641\u0642\u0637
2. \u0625\u0630\u0627 \u0644\u0645 \u062a\u062c\u062f \u0627\u0644\u0625\u062c\u0627\u0628\u0629\u060c \u0642\u0644 \u0630\u0644\u0643 \u0628\u0648\u0636\u0648\u062d
3. \u0627\u0633\u062a\u062e\u062f\u0645 \u0627\u0644\u0644\u063a\u0629 \u0627\u0644\u0639\u0631\u0628\u064a\u0629 \u0627\u0644\u0641\u0635\u062d\u0649
4. \u0643\u0646 \u0648\u0627\u0636\u062d\u0627\u064b \u0648\u0645\u0641\u0635\u0644\u0627\u064b
5. \u0627\u0633\u062a\u062e\u062f\u0645 \u0635\u064a\u063a\u0629 Markdown \u0641\u064a \u0627\u0644\u0625\u062c\u0627\u0628\u0629
{notes_section}

OUTPUT_FORMAT: Markdown
LANGUAGE: Arabic (\u0645\u0627 \u0644\u0645 \u064a\u064f\u062d\u062f\u062f \u062e\u0644\u0627\u0641 \u0630\u0644\u0643)

CONTEXT:
{context}

\u0627\u0644\u0633\u0624\u0627\u0644: {question}

\u0627\u0644\u0625\u062c\u0627\u0628\u0629:"""

        try:
            return self._generate_content(prompt, max_tokens=1000)
        except GeminiError as e:
            logger.error(f"Document Q&A failed: {e}")
            return "\u0639\u0630\u0631\u0627\u064b\u060c \u062d\u062f\u062b \u062e\u0637\u0623 \u0623\u062b\u0646\u0627\u0621 \u0645\u0639\u0627\u0644\u062c\u0629 \u0633\u0624\u0627\u0644\u0643. \u064a\u0631\u062c\u0649 \u0627\u0644\u0645\u062d\u0627\u0648\u0644\u0629 \u0645\u0631\u0629 \u0623\u062e\u0631\u0649."

    def _find_relevant_chunks(self, chunks: List[str], question: str) -> str:
        question_words = set(question.lower().split())
        scored = []
        for chunk in chunks:
            chunk_words = set(chunk.lower().split())
            overlap = len(question_words & chunk_words)
            scored.append((overlap, chunk))
        scored.sort(key=lambda x: x[0], reverse=True)
        top_chunks = [c[1] for c in scored[:3]]
        return "\n\n---\n\n".join(top_chunks)

    def generate_and_save_summary(self, file_obj, user_notes: str = "") -> AIResponse:
        """توليد ملخص وحفظه كملف .md."""
        text = self.extract_text_from_file(file_obj)
        if not text:
            return AIResponse(success=False, error='\u0644\u0627 \u064a\u0645\u0643\u0646 \u0627\u0633\u062a\u062e\u0631\u0627\u062c \u0627\u0644\u0646\u0635 \u0645\u0646 \u0627\u0644\u0645\u0644\u0641')

        try:
            summary = self.generate_summary(text, user_notes=user_notes)
            md_path = self._storage.save_summary(
                file_id=file_obj.id,
                content=summary,
                metadata={
                    'source_file': file_obj.title,
                    'course': str(file_obj.course),
                    'date': datetime.now().strftime('%Y-%m-%d %H:%M'),
                    'model': self._model_name,
                }
            )
            return AIResponse(success=True, data=summary, md_file_path=md_path)
        except Exception as e:
            logger.error(f"generate_and_save_summary failed: {e}")
            return AIResponse(success=False, error=str(e))

    def generate_and_save_questions(
        self, file_obj, matrix: QuestionMatrixConfig, user_notes: str = ""
    ) -> AIResponse:
        """توليد أسئلة وحفظها كملف .md."""
        text = self.extract_text_from_file(file_obj)
        if not text:
            return AIResponse(success=False, error='\u0644\u0627 \u064a\u0645\u0643\u0646 \u0627\u0633\u062a\u062e\u0631\u0627\u062c \u0627\u0644\u0646\u0635 \u0645\u0646 \u0627\u0644\u0645\u0644\u0641')

        try:
            questions = self.generate_questions_matrix(text, matrix, user_notes)
            if not questions:
                return AIResponse(success=False, error='\u0644\u0645 \u064a\u062a\u0645\u0643\u0646 AI \u0645\u0646 \u062a\u0648\u0644\u064a\u062f \u0623\u0633\u0626\u0644\u0629')

            md_path = self._storage.save_questions(
                file_id=file_obj.id,
                questions_data=questions,
                metadata={
                    'source_file': file_obj.title,
                    'course': str(file_obj.course),
                    'date': datetime.now().strftime('%Y-%m-%d %H:%M'),
                    'total_questions': str(len(questions)),
                    'total_score': str(matrix.total_score),
                    'model': self._model_name,
                }
            )
            return AIResponse(success=True, data=questions, md_file_path=md_path)
        except Exception as e:
            logger.error(f"generate_and_save_questions failed: {e}")
            return AIResponse(success=False, error=str(e))

    def test_connection(self) -> AIResponse:
        """اختبار الاتصال بـ Manus API Proxy."""
        try:
            start_ms = int(time.time() * 1000)
            response = self._generate_content("\u0642\u0644: \u0645\u0631\u062d\u0628\u0627\u064b\u060c \u0623\u0646\u0627 \u062c\u0627\u0647\u0632!", max_tokens=50)
            latency = int(time.time() * 1000) - start_ms
            return AIResponse(success=True, data={
                'response': response,
                'latency_ms': latency,
                'model': self._model_name,
                'base_url': MANUS_BASE_URL,
            })
        except GeminiError as e:
            return AIResponse(success=False, error=str(e))


# ========== Celery Tasks (Optional) ==========

try:
    from celery import shared_task
    CELERY_AVAILABLE = True
except ImportError:
    CELERY_AVAILABLE = False
    def shared_task(*args, **kwargs):
        def decorator(func):
            return func
        return decorator


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def generate_summary_async(self, file_id: int, user_notes: str = "") -> Dict[str, Any]:
    from apps.courses.models import LectureFile
    from apps.ai_features.models import AISummary

    try:
        file_obj = LectureFile.objects.get(pk=file_id)
        service = GeminiService()
        result = service.generate_and_save_summary(file_obj, user_notes=user_notes)

        if result.success:
            AISummary.objects.update_or_create(
                file=file_obj,
                defaults={
                    'summary_text': result.data[:200] + '...' if len(result.data) > 200 else result.data,
                    'is_cached': True,
                    'model_used': service._model_name,
                }
            )
            return {'success': True, 'md_file_path': result.md_file_path}
        return {'success': False, 'error': result.error}

    except LectureFile.DoesNotExist:
        return {'success': False, 'error': '\u0627\u0644\u0645\u0644\u0641 \u063a\u064a\u0631 \u0645\u0648\u062c\u0648\u062f'}
    except Exception as e:
        logger.error(f"Async summary generation failed: {e}")
        if CELERY_AVAILABLE and hasattr(self, 'retry'):
            raise self.retry(exc=e)
        return {'success': False, 'error': str(e)}


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def generate_questions_async(
    self, file_id: int, question_type: str = 'mixed',
    num_questions: int = 5, user_notes: str = ""
) -> Dict[str, Any]:
    from apps.courses.models import LectureFile
    from apps.ai_features.models import AIGeneratedQuestion

    try:
        file_obj = LectureFile.objects.get(pk=file_id)
        service = GeminiService()

        text = service.extract_text_from_file(file_obj)
        if not text:
            return {'success': False, 'error': '\u0644\u0627 \u064a\u0645\u0643\u0646 \u0627\u0633\u062a\u062e\u0631\u0627\u062c \u0627\u0644\u0646\u0635'}

        q_type = QuestionType(question_type) if question_type in [e.value for e in QuestionType] else QuestionType.MIXED
        questions = service.generate_questions(text, q_type, num_questions, user_notes)

        saved_ids = []
        for q in questions:
            ai_q = AIGeneratedQuestion.objects.create(
                file=file_obj,
                question_type=q.get('type', 'short_answer'),
                question_text=q.get('question', ''),
                options=q.get('options'),
                correct_answer=q.get('answer', ''),
                explanation=q.get('explanation', ''),
            )
            saved_ids.append(ai_q.id)

        return {'success': True, 'question_ids': saved_ids, 'count': len(saved_ids)}

    except LectureFile.DoesNotExist:
        return {'success': False, 'error': '\u0627\u0644\u0645\u0644\u0641 \u063a\u064a\u0631 \u0645\u0648\u062c\u0648\u062f'}
    except Exception as e:
        logger.error(f"Async question generation failed: {e}")
        if CELERY_AVAILABLE and hasattr(self, 'retry'):
            raise self.retry(exc=e)
        return {'success': False, 'error': str(e)}
